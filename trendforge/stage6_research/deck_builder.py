"""Build slides + speaker notes from a dossier."""
from __future__ import annotations

import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

log = logging.getLogger(__name__)


def slide_count_from_duration(duration_min: int | None) -> int:
    if not duration_min:
        return 12
    # ~1 slide per 90s, capped between 6 and 30
    n = max(6, min(30, int(duration_min * 60 / 90)))
    return n


def split_dossier_into_sections(dossier: str) -> list[tuple[str, str]]:
    """Return [(heading, body)] using top-level `##` headings."""
    parts = re.split(r"^##+\s+", dossier, flags=re.MULTILINE)
    if len(parts) <= 1:
        return [("Overview", dossier.strip())]
    sections: list[tuple[str, str]] = []
    # parts[0] is the preamble (under the H1)
    for chunk in parts[1:]:
        lines = chunk.strip().splitlines()
        if not lines:
            continue
        heading = lines[0].strip()
        body = "\n".join(lines[1:]).strip()
        sections.append((heading, body))
    return sections


def _add_text_slide(prs, title: str, body_lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:120]
    body = slide.placeholders[1].text_frame
    body.word_wrap = True
    body.text = body_lines[0][:300] if body_lines else ""
    for line in body_lines[1:]:
        p = body.add_paragraph()
        p.text = line[:300]
    # shrink font on overflow
    for para in body.paragraphs:
        for run in para.runs:
            run.font.size = Pt(16)


def build_slides_and_notes(dossier_path: Path, output_dir: Path,
                           topic: str, audience: str | None,
                           duration_min: int | None) -> tuple[Path, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    dossier = dossier_path.read_text(encoding="utf-8")
    sections = split_dossier_into_sections(dossier)
    target_count = slide_count_from_duration(duration_min)

    prs = Presentation()

    # Title
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_slide.placeholders[1].text = (
        f"For: {audience or 'technical audience'} · {duration_min or 30} min"
    )

    speaker_notes = [f"# Speaker notes — {topic}", ""]

    sections_to_use = sections[: max(target_count - 2, 4)]
    for idx, (heading, body) in enumerate(sections_to_use, start=1):
        bullets = [
            ln.lstrip("-* ").strip()
            for ln in body.splitlines()
            if ln.strip() and (ln.lstrip().startswith(("-", "*")) or len(ln) < 200)
        ]
        if not bullets:
            bullets = [body[:300]]
        bullets = bullets[:6]
        _add_text_slide(prs, heading, bullets)

        speaker_notes.append(f"## Slide {idx + 1}: {heading}")
        speaker_notes.append(body[:1200])
        speaker_notes.append("")

    # Closing slide
    _add_text_slide(prs, "Q&A", ["Open the floor.", "See qa_prep.md for predicted questions."])

    slides_path = output_dir / "slides.pptx"
    prs.save(str(slides_path))

    notes_path = output_dir / "speaker_notes.md"
    notes_path.write_text("\n".join(speaker_notes), encoding="utf-8")
    return slides_path, notes_path
